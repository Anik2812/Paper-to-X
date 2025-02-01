import PyPDF2
from pylatexenc.latex2text import LatexNodes2Text
import openai
from pptx import Presentation
from PIL import Image
import moviepy.editor as mp
import os
from dotenv import load_dotenv

load_dotenv()

# OpenAI configuration
openai.api_key = os.getenv("OPENAI_API_KEY")

async def process_pdf(file):
    try:
        pdf_reader = PyPDF2.PdfReader(file.file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text()
        return text
    except Exception as e:
        raise Exception(f"Error processing PDF: {str(e)}")

async def process_latex(file):
    try:
        content = await file.read()
        latex_text = content.decode("utf-8")
        text = LatexNodes2Text().latex_to_text(latex_text)
        return text
    except Exception as e:
        raise Exception(f"Error processing LaTeX: {str(e)}")

async def process_text(file):
    try:
        content = await file.read()
        return content.decode("utf-8")
    except Exception as e:
        raise Exception(f"Error processing text: {str(e)}")

async def generate_ppt(content, options):
    try:
        # Generate summary using OpenAI
        summary = await generate_summary(content, options.get("slides", 5))
        
        # Create PowerPoint
        prs = Presentation()
        for slide_content in summary:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_content["title"]
            slide.shapes.placeholders[1].text = slide_content["content"]
        
        # Save and return file path
        output_path = f"outputs/ppt/{os.urandom(16).hex()}.pptx"
        prs.save(output_path)
        return output_path
    except Exception as e:
        raise Exception(f"Error generating PPT: {str(e)}")

async def generate_podcast(content, options):
    try:
        # Generate script using OpenAI
        script = await generate_script(content, options.get("duration", 5))
        
        # Convert to audio (implementation depends on chosen TTS service)
        # For this example, we'll use a placeholder
        output_path = f"outputs/podcast/{os.urandom(16).hex()}.mp3"
        return output_path
    except Exception as e:
        raise Exception(f"Error generating podcast: {str(e)}")

async def generate_graphical(content, options):
    try:
        # Generate visual elements using OpenAI
        elements = await generate_visual_elements(content)
        
        # Create image
        img = Image.new('RGB', (800, 600), color='white')
        # Add visual elements to image
        
        output_path = f"outputs/graphical/{os.urandom(16).hex()}.png"
        img.save(output_path)
        return output_path
    except Exception as e:
        raise Exception(f"Error generating graphical abstract: {str(e)}")

async def generate_video(content, options):
    try:
        # Generate script and visuals
        script = await generate_script(content, options.get("duration", 3))
        
        # Create video using moviepy
        # This is a placeholder implementation
        output_path = f"outputs/video/{os.urandom(16).hex()}.mp4"
        return output_path
    except Exception as e:
        raise Exception(f"Error generating video: {str(e)}")

async def generate_summary(content, num_slides):
    try:
        response = await openai.ChatCompletion.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": "Create a presentation summary with clear slides."},
                {"role": "user", "content": f"Summarize this into {num_slides} slides:\n{content}"}
            ]
        )
        return response.choices[0].message.content
    except Exception as e:
        raise Exception(f"Error generating summary: {str(e)}")

async def generate_script(content, duration):
    try:
        response = await openai.ChatCompletion.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": f"Create a {duration}-minute script."},
                {"role": "user", "content": f"Convert this into a script:\n{content}"}
            ]
        )
        return response.choices[0].message.content
    except Exception as e:
        raise Exception(f"Error generating script: {str(e)}")

async def generate_visual_elements(content):
    try:
        response = await openai.ChatCompletion.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": "Create visual elements for a graphical abstract."},
                {"role": "user", "content": f"Generate visual elements for:\n{content}"}
            ]
        )
        return response.choices[0].message.content
    except Exception as e:
        raise Exception(f"Error generating visual elements: {str(e)}")